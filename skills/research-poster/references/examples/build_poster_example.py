#!/usr/bin/env python
"""
WORKED EXAMPLE for the research-poster skill: a self-contained, runnable demo that builds a
48x27 (16:9) scientific poster with NATIVE, editable PowerPoint charts. The study, authors,
institutions, and numbers are entirely fictional placeholders. Copy this and swap the DATA block,
CONTENT strings, and PALETTE tokens for a real poster. No external assets are required (the header
uses a text wordmark instead of a logo image; swap in add_picture(...) with a real logo when you have one).

Run:  python build_poster_example.py   ->  writes "example_poster.pptx"
Then: bash ../../scripts/render.sh example_poster.pptx example_poster.pdf example_poster.png 1.7
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# ---- DATA block (single source of truth; every chart/tile recomputes from here) ----
RESPONDED, TOTAL, REGIONS = 34, 120, 9
AVAIL = {"Offers telehealth": 28, "None": 6}                 # of RESPONDED (34)
MODALITY = {"Video": 17, "Phone": 11}                        # of those offering (28)
BARRIERS = [("Broadband access", 19), ("Reimbursement", 15), ("Staff training", 11)]  # of 34, overlap OK

# ---- TDD-lite: assert the displayed percentages recompute from the data ----
def pct(n, d): return round(n / d * 100)
assert pct(RESPONDED, TOTAL) == 28
assert pct(AVAIL["Offers telehealth"], RESPONDED) == 82 and pct(AVAIL["None"], RESPONDED) == 18
assert pct(MODALITY["Video"], sum(MODALITY.values())) == 61
print("numbers verified OK")

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "example_poster.pptx")

# ---- PALETTE + type (generic, professional; swap for the venue/author house style) ----
NAVY = RGBColor(0x1F, 0x3A, 0x5F); GOLD = RGBColor(0xC9, 0xA2, 0x4B); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A); FIELD = RGBColor(0xF3, 0xF2, 0xEE); CARD = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xB9, 0xB9, 0xB2); RULE = RGBColor(0xDD, 0xDC, 0xD5); MUTE = RGBColor(0x4A, 0x4A, 0x45)
FONT = "Arial"; BODY = 29; HEAD = 38

prs = Presentation(); prs.slide_width = Inches(48); prs.slide_height = Inches(27)
slide = prs.slides.add_slide(prs.slide_layouts[6]); shapes = slide.shapes

def rect(x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line if line else fill
    if line is None: sp.line.fill.background()
    sp.shadow.inherit = False; return sp

def textbox(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08); return tf

def setrun(r, text, size, bold=False, italic=False, color=INK):
    r.text = text; r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = FONT; r.font.color.rgb = color

def para(tf, runs, align=PP_ALIGN.LEFT, space_after=13, line=1.16, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line
    for a, kw in runs: setrun(p.add_run(), *a, **kw)
    return p
def R(t, s, **kw): return ((t, s), kw)

# ---- header + footer ----
rect(0, 0, 48, 27, FIELD)
HEADER_H = 4.2
rect(0, 0, 48, HEADER_H, NAVY); rect(0, HEADER_H, 48, 0.18, GOLD)
FOOT_Y = 25.7; rect(0, FOOT_Y, 48, 27 - FOOT_Y, NAVY); rect(0, FOOT_Y - 0.16, 48, 0.16, GOLD)
# text wordmark in place of a logo image (swap in add_picture for a real logo)
wm = textbox(38.5, 0.5, 8.6, HEADER_H - 1.0, MSO_ANCHOR.MIDDLE)
p = wm.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; setrun(p.add_run(), "UNIVERSITY", 34, bold=True, color=WHITE)
p2 = wm.add_paragraph(); p2.alignment = PP_ALIGN.RIGHT; setrun(p2.add_run(), "MEDICAL CENTER", 22, color=GOLD)
tf = textbox(0.95, 0.4, 37.0, HEADER_H - 0.6)
para(tf, [R("Telehealth Adoption in Rural Primary Care: A Clinic Director Survey", 46, bold=True, italic=True, color=WHITE)],
     space_after=11, line=0.98, first=True)
para(tf, [R("A. Author", 30, bold=True, color=WHITE), R("¹     ", 30, color=GOLD),
          R("B. Researcher", 30, bold=True, color=WHITE), R("²     ", 30, color=GOLD),
          R("C. Analyst", 30, bold=True, color=WHITE), R("¹", 30, color=GOLD)], space_after=8)
para(tf, [R("¹University Medical Center, Department of Family Medicine     "
            "²Regional Health Institute, Center for Rural Health", 22, color=RGBColor(0xEA, 0xE7, 0xDA))],
     space_after=0, line=1.12)

# ---- grid ----
CT_TOP = HEADER_H + 0.6; CT_BOT = FOOT_Y - 0.5
LM = 0.85; GUT = 0.6; CW = [12.55, 20.0, 12.55]
CX = [LM, LM + CW[0] + GUT, LM + CW[0] + CW[1] + 2 * GUT]
BAR_H = 1.30; HEAD_BLK = BAR_H + 0.10 + 0.32

def section_header(x, y, w, title):
    bar = rect(x, y, w, BAR_H, NAVY); tf = bar.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; setrun(p.add_run(), title, HEAD, bold=True, color=WHITE)
    rect(x, y + BAR_H, w, 0.10, GOLD)

def flow_column(x, w, blocks):
    total = sum(b[0] for b in blocks)
    gap = (CT_BOT - CT_TOP - total) / (len(blocks) - 1) if len(blocks) > 1 else 0
    y = CT_TOP
    for h, fn in blocks: fn(x, y, w, h); y += h + gap

def sec(title, build):
    def fn(x, y, w, h):
        section_header(x, y, w, title); build(textbox(x, y + HEAD_BLK, w, h - HEAD_BLK))
    return fn

def intro(tf):
    para(tf, [R("Telehealth expanded rapidly in primary care, yet uptake in rural clinics is uneven and "
                "poorly characterized. This survey quantifies telehealth availability, the dominant visit "
                "modality, and the barriers clinic directors report.", BODY)], align=PP_ALIGN.JUSTIFY, first=True)

def methods(tf):
    para(tf, [R("A voluntary survey was emailed to directors of all ", BODY), R("120", BODY, bold=True, color=NAVY),
              R(" rural primary-care clinics in the network, covering three domains: telehealth availability, "
                "primary visit modality, and adoption barriers. Responses were summarized as counts and "
                "proportions; data collection is ongoing and results are preliminary.", BODY)],
         align=PP_ALIGN.JUSTIFY, first=True)

def keyfind(tf):
    para(tf, [R(f"{RESPONDED} of {TOTAL} clinics ({pct(RESPONDED, TOTAL)}%)", BODY, bold=True, color=NAVY),
              R(f" responded, across {REGIONS} regions. Most offered telehealth "
                f"({pct(AVAIL['Offers telehealth'], RESPONDED)}%); video was the dominant modality "
                f"({pct(MODALITY['Video'], sum(MODALITY.values()))}%). Broadband access was the most common "
                "reported barrier. Adoption timelines ranged widely (see charts).", BODY)],
         align=PP_ALIGN.JUSTIFY, first=True)

def concl(tf):
    para(tf, [R("Bottom line:  ", BODY, bold=True, color=NAVY),
              R("Telehealth is widely available in these rural clinics, but infrastructure and reimbursement "
                "gaps persist. Targeted broadband and payment support may close the remaining adoption gap.",
                BODY, bold=True, color=INK)], align=PP_ALIGN.JUSTIFY, first=True)
    para(tf, [R("This preliminary analysis (28% response) is subject to response bias; further outreach is "
                "ongoing.", BODY)], align=PP_ALIGN.JUSTIFY)

REFS = [  # Demo placeholders; replace with verified, DOI-linked citations (see references/citations.md)
    "Author A, Coauthor B, et al. Representative title on rural telehealth adoption. J Example Med. 2024;12(3):100-110.",
    "Researcher C, Analyst D. A second representative reference. Example Health Policy. 2023;8(1):22-31.",
    "Regional Health Data Registry. Rural clinic telehealth dataset. Accessed 2026. example.org/registry.",
]
def refs(tf):
    for i, r in enumerate(REFS, 1):
        para(tf, [R(f"{i}. ", 22, bold=True, color=NAVY), R(r, 22, color=INK)], space_after=12, line=1.1, first=(i == 1))

# ---- results panel with native editable charts ----
def stat_tile(tx, ty, tw, th, big, small):
    rect(tx, ty, tw, th, CARD, line=RULE); tf = textbox(tx, ty, tw, th, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(8); setrun(p.add_run(), big, 66, bold=True, color=NAVY)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; setrun(p2.add_run(), small, 27, bold=True, color=INK)

def chart_caption(bx, by, bw, title, sub):
    tf = textbox(bx, by, bw, 0.82); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p.add_run(), title, 31, bold=True, color=INK)
    tf2 = textbox(bx, by + 0.80, bw, 0.58); p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p.add_run(), sub, 23, italic=True, color=MUTE)

def native_donut(bx, by, bw, bh, title, sub, cats, vals, colors, center):
    # % in the legend + big number in the hole; NO on-slice data labels, and NEVER set label position
    chart_caption(bx, by, bw, title, sub)
    chy = by + 1.5; chh = bh - 1.5
    cd = CategoryChartData(); cd.categories = cats; cd.add_series("s", vals)
    ch = shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(bx), Inches(chy), Inches(bw), Inches(chh), cd).chart
    ch.has_title = False; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False; ch.legend.font.size = Pt(24); ch.legend.font.name = FONT
    pl = ch.plots[0]; pl.vary_by_categories = False; pl.has_data_labels = False
    for i, c in enumerate(colors):
        pt = pl.series[0].points[i]; pt.format.fill.solid(); pt.format.fill.fore_color.rgb = c
    ct = textbox(bx, chy + (chh - 1.3) / 2 - 0.55, bw, 1.1, MSO_ANCHOR.MIDDLE)
    p = ct.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; setrun(p.add_run(), center, 58, bold=True, color=NAVY)

def native_bar(bx, by, bw, bh, title, sub, cats, vals, colors):
    chart_caption(bx, by, bw, title, sub)
    chy = by + 1.5; chh = bh - 1.5
    cd = CategoryChartData(); cd.categories = cats; cd.add_series("Share", vals)
    ch = shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(bx), Inches(chy), Inches(bw), Inches(chh), cd).chart
    ch.has_title = False; ch.has_legend = False
    sr = ch.plots[0].series[0]; sr.format.fill.solid(); sr.format.fill.fore_color.rgb = colors[0]
    for i, c in enumerate(colors): sr.points[i].format.fill.solid(); sr.points[i].format.fill.fore_color.rgb = c
    pl = ch.plots[0]; pl.has_data_labels = True; dl = pl.data_labels
    dl.show_value = True; dl.number_format = '0.0"%"'; dl.number_format_is_linked = False
    dl.font.size = Pt(30); dl.font.bold = True; dl.font.name = FONT; dl.font.color.rgb = NAVY
    ca = ch.category_axis; ca.tick_labels.font.size = Pt(25); ca.tick_labels.font.bold = True; ca.tick_labels.font.name = FONT
    va = ch.value_axis; va.tick_labels.font.size = Pt(21); va.tick_labels.font.name = FONT

def results_panel(x, y, w, h):
    section_header(x, y, w, "RESULTS")
    cy = y + HEAD_BLK; card_h = CT_BOT - cy
    rect(x, cy, w, card_h, CARD, line=RULE)
    pad = 0.75; ix = x + pad; iw = w - 2 * pad
    tiles_h, dr_h, br_h = 2.7, 7.4, 7.0
    g = (card_h - tiles_h - dr_h - br_h) / 3
    ty = cy + g; tg = 0.45; tw = (iw - 2 * tg) / 3
    stat_tile(ix, ty, tw, tiles_h, f"{RESPONDED} / {TOTAL}", f"clinics responding ({pct(RESPONDED, TOTAL)}%)")
    stat_tile(ix + tw + tg, ty, tw, tiles_h, str(REGIONS), "regions represented")
    stat_tile(ix + 2 * (tw + tg), ty, tw, tiles_h, "6 to 24 mo", "adoption timeline range")
    dy = ty + tiles_h + g; dw = iw / 2 - 0.2
    av = pct(AVAIL["Offers telehealth"], RESPONDED); md = pct(MODALITY["Video"], sum(MODALITY.values()))
    native_donut(ix, dy, dw, dr_h, "Telehealth Availability", f"n = {RESPONDED} responding clinics",
                 [f"Offers telehealth ({av}%)", f"None ({100 - av}%)"], list(AVAIL.values()), [NAVY, GRAY], f"{av}%")
    native_donut(ix + iw / 2 + 0.2, dy, dw, dr_h, "Primary Visit Modality", f"n = {sum(MODALITY.values())} offering clinics",
                 [f"Video ({md}%)", f"Phone ({100 - md}%)"], list(MODALITY.values()), [NAVY, GOLD], f"{md}%")
    bar_cats = [b[0] for b in BARRIERS]; bar_vals = [round(n / RESPONDED * 100, 1) for _, n in BARRIERS]
    native_bar(ix, dy + dr_h + g, iw, br_h, "Reported Barriers", f"share of responding clinics (n = {RESPONDED})",
               bar_cats, bar_vals, [NAVY, NAVY, GOLD])

# ---- assemble ----
flow_column(CX[0], CW[0], [(6.6, sec("INTRODUCTION", intro)), (7.0, sec("MATERIALS & METHODS", methods))])
flow_column(CX[1], CW[1], [(CT_BOT - CT_TOP, results_panel)])
flow_column(CX[2], CW[2], [(4.8, sec("KEY FINDINGS", keyfind)), (6.0, sec("CONCLUSIONS", concl)), (8.0, sec("REFERENCES", refs))])

tf = textbox(0.95, FOOT_Y, 46.1, 27 - FOOT_Y, MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
setrun(p.add_run(), "Example Conference 2026   ·   University Medical Center, Family Medicine   ·   "
       "Disclosures: none   ·   Preliminary data, collection ongoing", 24, color=WHITE)

prs.save(OUT); print("saved:", OUT, "| shapes:", len(slide.shapes))
